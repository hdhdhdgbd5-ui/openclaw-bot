"""
Export Service
Handles exporting visualizations and dashboards to various formats
"""

import os
from typing import Dict
from datetime import datetime

import plotly.io as pio
import kaleido
from pptx import Presentation
from pptx.util import Inches

from core.logging import logger

class ExportService:
    """Service for exporting visualizations and data"""
    
    def __init__(self):
        self.exports_dir = "exports"
        os.makedirs(self.exports_dir, exist_ok=True)
    
    async def export_dashboard_to_pdf(self, dashboard, include_raw_data: bool = False) -> str:
        """Export dashboard to PDF"""
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        
        filename = f"{self.exports_dir}/dashboard_{dashboard.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        
        c = canvas.Canvas(filename, pagesize=A4)
        width, height = A4
        
        # Title
        c.setFont("Helvetica-Bold", 24)
        c.drawString(72, height - 72, dashboard.name)
        
        # Description
        if dashboard.description:
            c.setFont("Helvetica", 12)
            c.drawString(72, height - 100, dashboard.description)
        
        # Widgets summary
        y_position = height - 150
        c.setFont("Helvetica-Bold", 14)
        c.drawString(72, y_position, f"Widgets: {len(dashboard.widgets)}")
        
        # Footer
        c.setFont("Helvetica", 10)
        c.drawString(72, 50, f"Exported from InsightGenius on {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        
        c.save()
        
        return filename
    
    async def export_dashboard_to_pptx(self, dashboard) -> str:
        """Export dashboard to PowerPoint"""
        filename = f"{self.exports_dir}/dashboard_{dashboard.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = dashboard.name
        slide.placeholders[1].text = dashboard.description or ""
        
        # Add slide for each widget
        for widget in dashboard.widgets:
            content_slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
            title_frame = title_box.text_frame
            title_frame.text = widget.title or f"Widget {widget.id}"
            
            # TODO: Add chart images when available
        
        prs.save(filename)
        
        return filename
    
    async def export_dataset_to_excel(self, dataset) -> str:
        """Export dataset to Excel"""
        import pandas as pd
        
        filename = f"{self.exports_dir}/dataset_{dataset.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        
        # Convert sample data to DataFrame
        if dataset.sample_data:
            df = pd.DataFrame(dataset.sample_data)
            df.to_excel(filename, index=False, engine='openpyxl')
        else:
            # Create empty file with headers
            df = pd.DataFrame(columns=[c.name for c in dataset.columns])
            df.to_excel(filename, index=False, engine='openpyxl')
        
        return filename
    
    async def export_dataset_to_csv(self, dataset) -> str:
        """Export dataset to CSV"""
        import pandas as pd
        
        filename = f"{self.exports_dir}/dataset_{dataset.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        
        if dataset.sample_data:
            df = pd.DataFrame(dataset.sample_data)
            df.to_csv(filename, index=False)
        else:
            df = pd.DataFrame(columns=[c.name for c in dataset.columns])
            df.to_csv(filename, index=False)
        
        return filename
    
    async def export_visualization_to_image(self, visualization, format: str, 
                                           width: int, height: int) -> str:
        """Export visualization as image"""
        filename = f"{self.exports_dir}/viz_{visualization.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{format}"
        
        # TODO: Generate actual visualization image using plotly/kaleido
        # For now, create placeholder
        
        if format == "svg":
            # Create simple SVG placeholder
            svg_content = f'''<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}">
                <rect width="100%" height="100%" fill="#f0f0f0"/>
                <text x="50%" y="50%" text-anchor="middle" font-size="20">{visualization.name}</text>
            </svg>'''
            with open(filename, 'w') as f:
                f.write(svg_content)
        else:
            # For PNG/JPG, we'd need to use PIL or similar
            from PIL import Image, ImageDraw, ImageFont
            
            img = Image.new('RGB', (width, height), color='#f0f0f0')
            draw = ImageDraw.Draw(img)
            
            # Add text
            text = visualization.name
            # Use default font
            draw.text((width/2 - 50, height/2), text, fill='#333333')
            
            img.save(filename, format.upper())
        
        return filename
